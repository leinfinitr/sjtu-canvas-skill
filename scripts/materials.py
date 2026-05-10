from __future__ import annotations

import subprocess
from dataclasses import dataclass
from pathlib import Path

SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".ppt", ".md", ".txt"}


@dataclass
class MaterialText:
    path: Path
    text: str


def _score_path(path: Path, query_terms: list[str]) -> int:
    haystack = f"{path.stem} {path.name}".lower()
    return sum(1 for term in query_terms if term.lower() in haystack)


def find_material_file(workspace: str | Path, query: str) -> Path:
    root = Path(workspace)
    terms = [term for term in query.split() if term]
    candidates = [p for p in root.rglob("*") if p.is_file() and p.suffix.lower() in SUPPORTED_EXTENSIONS]
    if not candidates:
        raise FileNotFoundError(f"No supported material files found under {root}")
    ranked = sorted(candidates, key=lambda p: (_score_path(p, terms), -len(p.name)), reverse=True)
    best = ranked[0]
    if terms and _score_path(best, terms) == 0:
        raise FileNotFoundError(f"No material matched query {query!r} under {root}")
    return best


def extract_material_text(path: str | Path) -> MaterialText:
    p = Path(path)
    suffix = p.suffix.lower()
    if suffix in {".md", ".txt"}:
        return MaterialText(path=p, text=p.read_text(encoding="utf-8", errors="replace"))
    if suffix == ".pdf":
        return MaterialText(path=p, text=_extract_pdf_text(p))
    if suffix in {".pptx", ".ppt"}:
        return MaterialText(path=p, text=_extract_ppt_text(p))
    raise ValueError(f"Unsupported material type: {p.suffix}")


def _extract_pdf_text(path: Path) -> str:
    try:
        import fitz  # PyMuPDF
    except ImportError as exc:
        raise RuntimeError("PDF extraction requires pymupdf. Install with: uv add pymupdf") from exc
    doc = fitz.open(path)
    chunks = []
    for idx, page in enumerate(doc, start=1):
        chunks.append(f"\n\n## Page {idx}\n" + page.get_text())
    return "".join(chunks).strip()


def _extract_ppt_text(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("PPTX extraction requires python-pptx. Install with: uv add python-pptx") from exc
    prs = Presentation(path)
    chunks = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        chunks.append(f"## Slide {idx}\n" + "\n".join(texts))
    return "\n\n".join(chunks).strip()
